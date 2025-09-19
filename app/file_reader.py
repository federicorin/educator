import os
import mimetypes
import zipfile
import rarfile
import py7zr
import fitz
import docx
import pytesseract
from PIL import Image
from pptx import Presentation
from moviepy import VideoFileClip, AudioFileClip, CompositeVideoClip
import speech_recognition as sr

def procesar_archivo(filepath):
    mime_type, _ = mimetypes.guess_type(filepath)
    ext = os.path.splitext(filepath)[1].lower()

    if ext == '.pdf':
        return extraer_pdf(filepath)
    elif ext in ['.doc', '.docx']:
        return extraer_word(filepath)
    elif ext == '.txt':
        return extraer_txt(filepath)
    elif ext in ['.jpg', '.jpeg', '.png']:
        return extraer_ocr(filepath)
    elif ext in ['.ppt', '.pptx']:
        return extraer_powerpoint(filepath)
    elif ext in ['.mp4', '.mov', '.avi']:
        return extraer_audio_video(filepath)
    elif ext in ['.mp3', '.wav', '.ogg']:
        return extraer_audio(filepath)
    elif ext in ['.zip', '.rar', '.7z']:
        return extraer_archivo_comprimido(filepath)
    else:
        return f"[Archivo {os.path.basename(filepath)} subido. No se pudo procesar su contenido automáticamente.]"


def extraer_pdf(filepath):
    text = ""
    with fitz.open(filepath) as doc:
        for page in doc:
            text += page.get_text()
    return text.strip()


def extraer_word(filepath):
    doc = docx.Document(filepath)
    return "\n".join([p.text for p in doc.paragraphs]).strip()


def extraer_txt(filepath):
    with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
        return f.read().strip()


def extraer_ocr(filepath):
    image = Image.open(filepath)
    return pytesseract.image_to_string(image)


def extraer_powerpoint(filepath):
    prs = Presentation(filepath)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text.strip()


def extraer_audio_video(filepath):
    video = VideoFileClip(filepath)
    audio_path = filepath + "_temp_audio.wav"
    video.audio.write_audiofile(audio_path, logger=None)
    texto = extraer_audio(audio_path)
    os.remove(audio_path)
    return texto


def extraer_audio(filepath):
    recognizer = sr.Recognizer()
    with sr.AudioFile(filepath) as source:
        audio = recognizer.record(source)
    try:
        return recognizer.recognize_google(audio, language="es-ES")
    except sr.UnknownValueError:
        return "[Audio no reconocible.]"
    except sr.RequestError:
        return "[Error al conectar con el servicio de reconocimiento de voz.]"


def extraer_archivo_comprimido(filepath):
    base = os.path.splitext(filepath)[0] + "_extraido"
    os.makedirs(base, exist_ok=True)

    try:
        if filepath.endswith('.zip'):
            with zipfile.ZipFile(filepath, 'r') as z:
                z.extractall(base)
        elif filepath.endswith('.rar'):
            with rarfile.RarFile(filepath, 'r') as r:
                r.extractall(base)
        elif filepath.endswith('.7z'):
            with py7zr.SevenZipFile(filepath, 'r') as z:
                z.extractall(path=base)
    except Exception:
        return "[No se pudo descomprimir el archivo.]"

    texto = ""
    for root, _, files in os.walk(base):
        for file in files:
            ruta = os.path.join(root, file)
            try:
                texto += f"\n--- {file} ---\n"
                texto += procesar_archivo(ruta) + "\n"
            except Exception:
                texto += f"\n[No se pudo procesar {file}]\n"
    return texto.strip()
